import os
import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

def create_mock_assets():
    os.makedirs("test_assets", exist_ok=True)
    
    # 1. Create mock PPTX template
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI PowerPoint Builder Demo"
    subtitle.text = "Template Presentation"
    
    # Slide 2: Executive Summary (Slide Layout 1: Title and Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1]
    body.text = "{{bullets}}"
    
    # Slide 3: Revenue Analysis
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Revenue Analysis"
    body = slide.placeholders[1]
    body.text = "Detailed revenue overview"
    
    # Slide 4: Regional Analysis
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Regional Analysis"
    body = slide.placeholders[1]
    body.text = "Regional performance summary"

    # Slide 5: Trends
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Financial Trends"
    body = slide.placeholders[1]
    body.text = "Trends over time"

    # Slide 6: Recommendations
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Recommendations"
    body = slide.placeholders[1]
    body.text = "{{recommendations}}"
    
    prs.save("test_assets/template.pptx")
    print("Created test_assets/template.pptx")
    
    # 2. Create mock Excel data
    data = {
        "Date": [
            "2026-01-01", "2026-02-01", "2026-03-01", "2026-04-01", "2026-05-01", "2026-06-01",
            "2026-01-01", "2026-02-01", "2026-03-01", "2026-04-01", "2026-05-01", "2026-06-01"
        ],
        "Region": [
            "North", "North", "North", "North", "North", "North",
            "South", "South", "South", "South", "South", "South"
        ],
        "Sales": [
            12000, 15000, 14000, 18000, 22000, 25000,
            8000, 9500, 11000, 10500, 13000, 15000
        ],
        "Profit": [
            2400, 3100, 2900, 3800, 4800, 5600,
            1200, 1500, 1800, 1700, 2200, 2600
        ],
        "Margin_Pct": [
            20.0, 20.7, 20.7, 21.1, 21.8, 22.4,
            15.0, 15.8, 16.4, 16.2, 16.9, 17.3
        ]
    }
    
    df = pd.DataFrame(data)
    df.to_excel("test_assets/data.xlsx", index=False)
    print("Created test_assets/data.xlsx")

if __name__ == "__main__":
    create_mock_assets()
