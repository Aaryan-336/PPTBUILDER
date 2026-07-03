import os
import logging
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
from typing import List, Dict, Any
from pptx.chart.data import CategoryChartData
from .chart_generator import add_native_chart, create_chart_data

DUMMY_TEXT_INDICATORS = [
    "lorem ipsum", "dolor sit", "consectetur adipiscing", "dummy text",
    "placeholder", "click to add", "double click", "add text",
    "enter title", "enter subtitle", "insert text", "your text here",
    "add text here", "description here", "text box", "sample text",
    "your caption", "heading goes here", "subheading goes here",
    "[insert", "your description"
]

def is_dummy_text(text: str) -> bool:
    if not text:
        return False
    text_lower = text.lower().strip()
    return any(indicator in text_lower for indicator in DUMMY_TEXT_INDICATORS)

def check_overlap(rect1, rect2) -> bool:
    """Checks if two rectangles overlap. rect = (left, top, width, height)"""
    l1, t1, w1, h1 = rect1
    l2, t2, w2, h2 = rect2
    return not (l1 >= l2 + w2 or l2 >= l1 + w1 or t1 >= t2 + h2 or t2 >= t1 + h1)

def is_content_shape(shape, written_shape_ids) -> bool:
    """Checks if a shape contains actual content that should not be overlapped."""
    if shape.shape_id in written_shape_ids:
        return True
    if shape.has_text_frame:
        txt = shape.text.strip()
        if txt and not is_dummy_text(txt):
            return True
    if shape.has_table or shape.has_chart:
        return True
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    return False

logger = logging.getLogger("ppt_generator")

def get_text_frame_formatting(tf) -> Dict[str, Any]:
    """Retrieves font name, size, bold, italic, alignment, and color from the first run of the first paragraph of a text frame."""
    formatting = {
        "font_name": None,
        "font_size": None,
        "bold": None,
        "italic": None,
        "rgb": None,
        "theme_color": None,
        "alignment": None
    }
    try:
        if tf.paragraphs:
            p = tf.paragraphs[0]
            formatting["alignment"] = p.alignment
            # Try to get from first run if available, otherwise paragraph level
            if p.runs:
                font = p.runs[0].font
            else:
                font = p.font
                
            formatting["font_name"] = font.name
            formatting["font_size"] = font.size
            formatting["bold"] = font.bold
            formatting["italic"] = font.italic
            if font.color:
                try:
                    from pptx.enum.dml import MSO_COLOR_TYPE
                    if font.color.type == MSO_COLOR_TYPE.RGB:
                        formatting["rgb"] = font.color.rgb
                    elif font.color.type == MSO_COLOR_TYPE.SCHEME:
                        formatting["theme_color"] = font.color.theme_color
                except Exception:
                    pass
    except Exception as e:
        logger.warning(f"Failed to extract text frame formatting: {e}")
    return formatting

def apply_formatting_to_run(run, formatting: Dict[str, Any]):
    """Applies captured formatting to a run."""
    try:
        if formatting.get("font_name"):
            run.font.name = formatting["font_name"]
        if formatting.get("font_size"):
            run.font.size = formatting["font_size"]
        if formatting.get("bold") is not None:
            run.font.bold = formatting["bold"]
        if formatting.get("italic") is not None:
            run.font.italic = formatting["italic"]
        if formatting.get("rgb") is not None:
            run.font.color.rgb = formatting["rgb"]
        elif formatting.get("theme_color") is not None:
            run.font.color.theme_color = formatting["theme_color"]
    except Exception as e:
        logger.warning(f"Failed to apply formatting to run: {e}")

def adjust_font_size(original_size, text_len: int, default_size_pt: int):
    """Scales down font size by up to 20% (minimum 0.8 scale) and enforces minimum 14pt."""
    from pptx.util import Pt
    base_size = original_size if original_size is not None else Pt(default_size_pt)
    
    try:
        base_pt = base_size.pt if hasattr(base_size, 'pt') else (base_size / 12700)
    except Exception:
        base_pt = default_size_pt
        
    scale = 1.0
    if default_size_pt >= 28:  # Title-like
        if text_len > 40:
            scale = 0.8
        elif text_len > 20:
            scale = 0.9
    else:  # Body-like
        if text_len > 100:
            scale = 0.8
        elif text_len > 50:
            scale = 0.9
            
    new_pt = base_pt * scale
    
    # Enforce minimum font size of 14pt
    if new_pt < 14:
        new_pt = 14
        
    # Enforce never increasing size
    if new_pt > base_pt:
        new_pt = base_pt
        
    return Pt(round(new_pt))

def find_placeholder_by_type(slide, ph_type) -> Any:
    """Finds a slide placeholder by type (e.g., PP_PLACEHOLDER.TITLE)."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type == ph_type:
                    return shape
            except Exception:
                pass
    return None

def find_placeholder_by_index(slide, index: int) -> Any:
    """Finds a slide placeholder by its layout index."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.idx == index:
                    return shape
            except Exception:
                pass
    return None

def find_custom_title_shape(slide, original_title: str) -> Any:
    """Finds a slide title shape when no placeholder title exists."""
    if original_title:
        orig_lower = original_title.lower().strip()
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.lower().strip() == orig_lower:
                return shape
    # Fallback: find the top-most shape with text
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if text_shapes:
        return min(text_shapes, key=lambda s: s.top)
    return None

def find_custom_content_shapes(slide, title_shape) -> List[Any]:
    """Finds content text boxes on slides without placeholders."""
    content_shapes = []
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if not shape.has_text_frame:
            continue
        txt = shape.text.strip()
        if not txt:
            continue
        # Skip small helper tags, numbers, design instructions, colors
        if len(txt) <= 2 and txt.isdigit():
            continue
        if txt.lower() in ("credits", "fonts", "resource page", "design elements", "colors"):
            continue
        content_shapes.append(shape)
    # Sort from left-to-right, then top-to-bottom
    return sorted(content_shapes, key=lambda s: (s.top, s.left))

def generate_presentation(
    template_path: str,
    output_path: str,
    dataset_path: str,
    slides_config: List[Dict[str, Any]]
) -> str:
    """
    Reads the PowerPoint template, populates content and charts,
    and saves the finalized presentation.
    """
    logger.info(f"Generating presentation from {template_path} to {output_path}")

    # Load dataset
    try:
        if dataset_path.endswith(".csv"):
            df = pd.read_csv(dataset_path)
        else:
            df = pd.read_excel(dataset_path)
    except Exception as e:
        logger.error(f"Failed to read dataset: {e}")
        raise ValueError(f"Could not load dataset for presentation: {e}")

    # Load template presentation
    try:
        prs = Presentation(template_path)
    except Exception as e:
        logger.error(f"Failed to load presentation template: {e}")
        raise ValueError(f"Could not load presentation template: {e}")

    # Sort slides_config by slide_index to match slide numbers
    slides_config = sorted(slides_config, key=lambda x: x["slide_index"])

    for s_config in slides_config:
        slide_idx = s_config["slide_index"] - 1  # 0-indexed
        
        # Check if slide index is valid in template
        if slide_idx >= len(prs.slides):
            logger.warning(f"Config references slide index {slide_idx + 1}, but template only has {len(prs.slides)} slides. Skipping.")
            continue
            
        slide = prs.slides[slide_idx]
        if s_config.get("skip") or s_config.get("locked"):
            logger.info(f"Slide {s_config['slide_index']} is flagged as 'do not change'. Skipping modifications.")
            continue
        insights = s_config.get("insights", {})
        chart_config = s_config.get("chart")

        # Track written shape IDs to avoid deleting them later during cleanup
        written_shape_ids = set()

        # 1. Update Title
        title_text = insights.get("title")
        if title_text:
            # Try setting native title placeholder
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    tf = slide.shapes.title.text_frame
                    fmt = get_text_frame_formatting(tf)
                    fmt["font_size"] = adjust_font_size(fmt["font_size"], len(title_text), 36)
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = title_text
                    if p.runs:
                        apply_formatting_to_run(p.runs[0], fmt)
                    if fmt.get("alignment") is not None:
                        p.alignment = fmt["alignment"]
                    written_shape_ids.add(slide.shapes.title.shape_id)
            except Exception:
                # Fallback: search for first shape matching title format
                title_ph = find_placeholder_by_type(slide, PP_PLACEHOLDER.TITLE)
                if title_ph and title_ph.has_text_frame:
                    tf = title_ph.text_frame
                    fmt = get_text_frame_formatting(tf)
                    fmt["font_size"] = adjust_font_size(fmt["font_size"], len(title_text), 36)
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = title_text
                    if p.runs:
                        apply_formatting_to_run(p.runs[0], fmt)
                    if fmt.get("alignment") is not None:
                        p.alignment = fmt["alignment"]
                    written_shape_ids.add(title_ph.shape_id)

        # 2. Update Bullets and Recommendations
        bullets = insights.get("bullets", [])
        recommendations = insights.get("recommendations", [])

        # Find body text placeholders
        body_phs = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    p_type = shape.placeholder_format.type
                    if p_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.OBJECT) and shape.has_text_frame:
                        body_phs.append(shape)
                except Exception:
                    pass

        # Sort body placeholders by index or layout top coordinate
        body_phs = sorted(body_phs, key=lambda s: (s.top, s.left))

        # Insert Bullets into the primary body placeholder
        if bullets and body_phs:
            primary_ph = body_phs[0]
            tf = primary_ph.text_frame
            fmt = get_text_frame_formatting(tf)
            
            # Apply dynamic font scaling based on total bullets text length
            total_bullets_len = sum(len(b) for b in bullets)
            fmt["font_size"] = adjust_font_size(fmt["font_size"], total_bullets_len, 16)
            
            tf.clear()  # removes default template placeholder text
            
            for idx, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                if p.runs:
                    apply_formatting_to_run(p.runs[0], fmt)
                if fmt.get("alignment") is not None:
                    p.alignment = fmt["alignment"]
            written_shape_ids.add(primary_ph.shape_id)
                
            # If recommendations exist and there is a second placeholder, populate it
            if recommendations and len(body_phs) > 1:
                secondary_ph = body_phs[1]
                tf_sec = secondary_ph.text_frame
                fmt_sec = get_text_frame_formatting(tf_sec)
                
                total_recs_len = sum(len(r) for r in recommendations)
                fmt_sec["font_size"] = adjust_font_size(fmt_sec["font_size"], total_recs_len, 16)
                
                tf_sec.clear()
                for idx, rec in enumerate(recommendations):
                    p = tf_sec.paragraphs[0] if idx == 0 else tf_sec.add_paragraph()
                    p.text = rec
                    p.level = 0
                    if p.runs:
                        apply_formatting_to_run(p.runs[0], fmt_sec)
                    if fmt_sec.get("alignment") is not None:
                        p.alignment = fmt_sec["alignment"]
                written_shape_ids.add(secondary_ph.shape_id)
            elif recommendations:
                # Fallback: Append recommendations to primary placeholder if no secondary is present
                p_space = tf.add_paragraph()
                p_space.text = ""
                p_head = tf.add_paragraph()
                p_head.text = "Recommendations:"
                p_head.level = 0
                if p_head.runs:
                    apply_formatting_to_run(p_head.runs[0], fmt)
                    p_head.runs[0].font.bold = True
                for rec in recommendations:
                    p = tf.add_paragraph()
                    p.text = rec
                    p.level = 1
                    if p.runs:
                        apply_formatting_to_run(p.runs[0], fmt)
                    if fmt.get("alignment") is not None:
                        p.alignment = fmt["alignment"]
        else:
            # Custom shape fallback when native layout placeholders are missing
            title_ph = find_custom_title_shape(slide, s_config.get("slide_title"))
            if title_ph and title_text:
                tf = title_ph.text_frame
                fmt = get_text_frame_formatting(tf)
                fmt["font_size"] = adjust_font_size(fmt["font_size"], len(title_text), 36)
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_text
                if p.runs:
                    apply_formatting_to_run(p.runs[0], fmt)
                if fmt.get("alignment") is not None:
                    p.alignment = fmt["alignment"]
                written_shape_ids.add(title_ph.shape_id)
                
            content_shapes = find_custom_content_shapes(slide, title_ph)
            if content_shapes:
                if len(content_shapes) == 1:
                    shape = content_shapes[0]
                    tf = shape.text_frame
                    fmt = get_text_frame_formatting(tf)
                    
                    total_text_len = sum(len(b) for b in bullets) + sum(len(r) for r in recommendations)
                    fmt["font_size"] = adjust_font_size(fmt["font_size"], total_text_len, 16)
                    
                    tf.clear()
                    for idx, b in enumerate(bullets):
                        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                        p.text = b
                        if p.runs:
                            apply_formatting_to_run(p.runs[0], fmt)
                        if fmt.get("alignment") is not None:
                            p.alignment = fmt["alignment"]
                    if recommendations:
                        p_space = tf.add_paragraph()
                        p_space.text = ""
                        p_head = tf.add_paragraph()
                        p_head.text = "Recommendations:"
                        if p_head.runs:
                            apply_formatting_to_run(p_head.runs[0], fmt)
                            p_head.runs[0].font.bold = True
                        for rec in recommendations:
                            p = tf.add_paragraph()
                            p.text = rec
                            if p.runs:
                                apply_formatting_to_run(p.runs[0], fmt)
                            if fmt.get("alignment") is not None:
                                p.alignment = fmt["alignment"]
                    written_shape_ids.add(shape.shape_id)
                elif len(content_shapes) == 2:
                    shape1 = content_shapes[0]
                    shape2 = content_shapes[1]
                    # Shape 1 gets bullets
                    tf1 = shape1.text_frame
                    fmt1 = get_text_frame_formatting(tf1)
                    
                    total_bullets_len = sum(len(b) for b in bullets)
                    fmt1["font_size"] = adjust_font_size(fmt1["font_size"], total_bullets_len, 16)
                    
                    tf1.clear()
                    for idx, b in enumerate(bullets):
                        p = tf1.paragraphs[0] if idx == 0 else tf1.add_paragraph()
                        p.text = b
                        if p.runs:
                            apply_formatting_to_run(p.runs[0], fmt1)
                        if fmt1.get("alignment") is not None:
                            p.alignment = fmt1["alignment"]
                    written_shape_ids.add(shape1.shape_id)
                    # Shape 2 gets recommendations
                    tf2 = shape2.text_frame
                    fmt2 = get_text_frame_formatting(tf2)
                    
                    total_recs_len = sum(len(r) for r in recommendations)
                    fmt2["font_size"] = adjust_font_size(fmt2["font_size"], total_recs_len, 16)
                    
                    tf2.clear()
                    for idx, r in enumerate(recommendations):
                        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
                        p.text = r
                        if p.runs:
                            apply_formatting_to_run(p.runs[0], fmt2)
                        if fmt2.get("alignment") is not None:
                            p.alignment = fmt2["alignment"]
                    written_shape_ids.add(shape2.shape_id)
                else:
                    # Distribute bullets + recommendations
                    all_items = bullets + recommendations
                    for idx, shape in enumerate(content_shapes):
                        tf = shape.text_frame
                        fmt = get_text_frame_formatting(tf)
                        
                        if idx < len(all_items):
                            item_text = all_items[idx]
                            fmt["font_size"] = adjust_font_size(fmt["font_size"], len(item_text), 16)
                            tf.clear()
                            p = tf.paragraphs[0]
                            p.text = item_text
                            if p.runs:
                                apply_formatting_to_run(p.runs[0], fmt)
                            if fmt.get("alignment") is not None:
                                p.alignment = fmt["alignment"]
                            written_shape_ids.add(shape.shape_id)
                        else:
                            tf.clear()

        # 3. Tag-based template replacement (supports {{variable}} in text frames)
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                replaced = False
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if "{{title}}" in run.text.lower() and title_text:
                            run.text = run.text.replace("{{title}}", title_text)
                            replaced = True
                        if "{{bullets}}" in run.text.lower() and bullets:
                            run.text = run.text.replace("{{bullets}}", "\n".join(bullets))
                            replaced = True
                        if "{{recommendations}}" in run.text.lower() and recommendations:
                            run.text = run.text.replace("{{recommendations}}", "\n".join(recommendations))
                            replaced = True
                if replaced:
                    written_shape_ids.add(shape.shape_id)

        # 4. Generate & Insert Chart
        chart_placed = False
        chart_rect = None

        if chart_config and chart_config != "skip":
            # Check if there is an existing native chart shape on the slide to update in-place
            existing_chart_shape = None
            for shape in slide.shapes:
                if shape.has_chart:
                    existing_chart_shape = shape
                    break
            
            if existing_chart_shape:
                try:
                    chart = existing_chart_shape.chart
                    x_col = chart_config.get("x_axis")
                    y_col = chart_config.get("y_axis")
                    aggregation = chart_config.get("aggregation") or "sum"
                    categories, values = create_chart_data(df, x_col, y_col, aggregation)
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = categories
                    chart_data.add_series(f"{y_col} ({aggregation.title()})", tuple(values))
                    
                    chart.replace_data(chart_data)
                    written_shape_ids.add(existing_chart_shape.shape_id)
                    chart_placed = True
                    logger.info("Successfully updated existing chart in-place.")
                except Exception as chart_err:
                    logger.warning(f"Failed to update existing chart in-place: {chart_err}. Falling back to creation.")

            if not chart_placed:
                # Check for a chart placeholder
                chart_ph = find_placeholder_by_type(slide, PP_PLACEHOLDER.CHART)
                
                # If no chart placeholder, look for standard body or object placeholders
                if not chart_ph and len(body_phs) > 1:
                    # If slide has multiple placeholders, use the last one for the chart
                    # But make sure we didn't write bullets to it!
                    candidate = body_phs[-1]
                    if candidate.shape_id not in written_shape_ids:
                        chart_ph = candidate
                    
                if chart_ph:
                    left = chart_ph.left
                    top = chart_ph.top
                    width = chart_ph.width
                    height = chart_ph.height
                    chart_rect = (left, top, width, height)
                else:
                    # Absolute positioning fallback (standard RHS or lower slide half)
                    left = Inches(5.0)
                    top = Inches(1.5)
                    width = Inches(4.5)
                    height = Inches(3.5)
                    chart_rect = (left, top, width, height)

                # Check if this chart overlaps with any existing layout content
                # Rule: Never place charts over existing content (Rule 7)
                overlap_found = False
                for shape in slide.shapes:
                    if chart_ph and shape.shape_id == chart_ph.shape_id:
                        continue
                    try:
                        if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                            continue
                    except Exception:
                        pass
                    if is_content_shape(shape, written_shape_ids):
                        try:
                            shape_rect = (shape.left, shape.top, shape.width, shape.height)
                            if check_overlap(shape_rect, chart_rect):
                                overlap_found = True
                                break
                        except Exception:
                            pass

                if overlap_found:
                    logger.warning(f"Chart insertion skipped for slide {s_config['slide_index']}: placement would overlap existing layout content.")
                else:
                    # If placeholder existed, delete it first
                    if chart_ph:
                        try:
                            sp_tree = slide.shapes._spTree
                            sp_tree.remove(chart_ph._element)
                        except Exception as del_err:
                            logger.warning(f"Could not delete placeholder shape for chart placement: {del_err}")

                    # Draw the chart
                    chart_shape = add_native_chart(slide, df, chart_config, left, top, width, height)
                    if chart_shape:
                        written_shape_ids.add(chart_shape.shape_id)
                        chart_placed = True

        # 6. Dummy Content & Unpopulated Placeholder Cleanup
        to_delete = set()
        for shape in slide.shapes:
            if shape.shape_id in written_shape_ids:
                continue
            
            # Check if shape is a body placeholder that was not written to
            if shape in body_phs:
                to_delete.add(shape)
                continue
                
            # Check if text shape contains dummy text
            if shape.has_text_frame:
                txt = shape.text.strip()
                if txt and is_dummy_text(txt):
                    to_delete.add(shape)
                elif not txt and shape.is_placeholder:
                    # Clear unpopulated placeholders of any type (e.g. subtitle, footer)
                    to_delete.add(shape)
            
            # Also check if it's a chart shape we wanted to replace but couldn't, or generic overlapping unwritten text
            if chart_placed and chart_rect and shape.has_text_frame and not shape.is_placeholder:
                # If there's an unwritten text box directly overlapping the chart, delete it
                try:
                    shape_rect = (shape.left, shape.top, shape.width, shape.height)
                    if check_overlap(shape_rect, chart_rect):
                        to_delete.add(shape)
                except Exception:
                    pass

        # Execute deletion cleanly
        for shape in to_delete:
            try:
                sp_tree = slide.shapes._spTree
                sp_tree.remove(shape._element)
                logger.info(f"Cleaned up redundant shape '{shape.name}' from slide.")
            except Exception as del_err:
                logger.warning(f"Could not delete shape '{shape.name}': {del_err}")

    # Save presentation
    try:
        prs.save(output_path)
    except Exception as e:
        logger.error(f"Failed to save final presentation: {e}")
        raise ValueError(f"Could not write output presentation file: {e}")

    return output_path
