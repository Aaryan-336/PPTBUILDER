import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Optional

logger = logging.getLogger("ppt_parser")

def get_placeholder_type_name(placeholder_format_type) -> str:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for name, value in PP_PLACEHOLDER.__members__.items():
            if value == placeholder_format_type:
                return name
        return str(placeholder_format_type)
    except Exception:
        return str(placeholder_format_type)

def parse_ppt_file(file_path: str) -> Dict[str, Any]:
    """
    Parses a PPTX file and returns a structured dictionary of its slides,
    titles, layouts, placeholders, and shapes.
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.error(f"Failed to open PPTX file {file_path}: {e}")
        raise ValueError(f"Invalid PowerPoint file: {e}")

    parsed_slides = []
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i + 1,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "title": None,
            "placeholders": [],
            "other_shapes": []
        }
        
        # Determine title
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_info["title"] = slide.shapes.title.text.strip()
        except Exception:
            pass  # Some slides don't have a title placeholder
            
        for shape in slide.shapes:
            shape_text = None
            if shape.has_text_frame:
                shape_text = shape.text.strip()
                
            # If we couldn't get title from shapes.title, fallback to looking at PP_PLACEHOLDER.TITLE
            if not slide_info["title"] and shape.is_placeholder:
                try:
                    p_type = shape.placeholder_format.type
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    if p_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and shape_text:
                        slide_info["title"] = shape_text
                except Exception:
                    pass

            shape_info = {
                "index": getattr(shape, "shape_id", -1),
                "name": shape.name,
                "type": str(shape.shape_type),
                "text": shape_text,
                "has_text": shape.has_text_frame
            }
            
            # Layout coordinates
            try:
                shape_info["x"] = shape.left
                shape_info["y"] = shape.top
                shape_info["width"] = shape.width
                shape_info["height"] = shape.height
            except Exception:
                pass

            if shape.is_placeholder:
                try:
                    ph_format = shape.placeholder_format
                    shape_info["placeholder_idx"] = ph_format.idx
                    shape_info["placeholder_type"] = get_placeholder_type_name(ph_format.type)
                    slide_info["placeholders"].append(shape_info)
                except Exception as ph_err:
                    logger.warning(f"Error reading placeholder format for shape {shape.name}: {ph_err}")
                    slide_info["other_shapes"].append(shape_info)
            else:
                slide_info["other_shapes"].append(shape_info)
                
        parsed_slides.append(slide_info)
        
    return {"slides": parsed_slides}
